#!/usr/bin/env python3
import argparse
import json
import sys
from pathlib import Path
import zipfile

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError as exc:
    print("python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    raise


EMU_PER_PT = 12700.0
ANALYSIS_VERSION = 2
MAX_BODY_COLUMNS = 4


def emu_to_pt(value):
    try:
        return float(value) / EMU_PER_PT
    except Exception:
        return 0.0


def shape_area(shape):
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def pick_font_pt(shape, default_pt=18.0):
    if not getattr(shape, "has_text_frame", False):
        return default_pt
    try:
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return default_pt
        paragraph = text_frame.paragraphs[0]
        if paragraph.runs:
            size = paragraph.runs[0].font.size
            if size is not None:
                return size.pt
        if paragraph.font and paragraph.font.size is not None:
            return paragraph.font.size.pt
    except Exception:
        return default_pt
    return default_pt


def estimate_capacity(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    font_pt = pick_font_pt(shape)
    width_pt = emu_to_pt(getattr(shape, "width", 0))
    height_pt = emu_to_pt(getattr(shape, "height", 0))
    tf = shape.text_frame
    margin_left = emu_to_pt(getattr(tf, "margin_left", 0))
    margin_right = emu_to_pt(getattr(tf, "margin_right", 0))
    margin_top = emu_to_pt(getattr(tf, "margin_top", 0))
    margin_bottom = emu_to_pt(getattr(tf, "margin_bottom", 0))
    usable_width = max(width_pt - margin_left - margin_right, font_pt)
    usable_height = max(height_pt - margin_top - margin_bottom, font_pt * 1.2)
    line_height = max(font_pt * 1.2, 10.0)
    max_lines = max(1, int(usable_height / line_height))
    char_width = max(font_pt * 0.9, 8.0)
    max_chars = max(4, int(usable_width / char_width))
    return {
        "lines": max_lines,
        "chars": max_chars,
    }


def placeholder_type_name(ph_type):
    try:
        return PP_PLACEHOLDER(ph_type).name
    except Exception:
        return str(ph_type)


def collect_placeholder_shapes(shapes):
    placeholders = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        placeholders.append(shape)
    return placeholders


def collect_text_shapes(shapes, exclude=None):
    exclude = list(exclude) if exclude else []
    result = []
    for shape in shapes:
        if any(shape is item for item in exclude):
            continue
        if getattr(shape, "has_text_frame", False):
            result.append(shape)
    return result


def classify_placeholders(placeholders, title_shapes):
    body_types = []
    for name in ("BODY", "CONTENT", "TEXT", "OBJECT"):
        if hasattr(PP_PLACEHOLDER, name):
            body_types.append(getattr(PP_PLACEHOLDER, name))
    body_types = tuple(body_types)

    image_types = []
    for name in ("PICTURE", "CHART", "TABLE", "CLIP_ART", "MEDIA"):
        if hasattr(PP_PLACEHOLDER, name):
            image_types.append(getattr(PP_PLACEHOLDER, name))
    image_types = tuple(image_types)

    titles = []
    bodies = []
    has_image = False

    for shape in placeholders:
        if any(shape is item for item in title_shapes):
            titles.append(shape)
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in body_types:
            bodies.append(shape)
        elif ph_type in image_types:
            has_image = True
        else:
            # treat unknown placeholders with text frame as body
            if getattr(shape, "has_text_frame", False):
                bodies.append(shape)

    return titles, bodies, has_image


def should_keep_body_entry(entry, canvas_area):
    lines = entry.get("lines", 0)
    chars = entry.get("chars", 0)
    area = entry.get("area", 0)
    ratio = (float(area) / float(canvas_area)) if canvas_area else 0.0
    is_placeholder = bool(entry.get("is_placeholder", False))

    # Tiny one-line shapes are usually decorative labels or numbering fragments.
    if lines <= 1 and chars <= 12:
        return False

    # Non-placeholder tiny text boxes are often ornaments in designer templates.
    if not is_placeholder and ratio < 0.01 and (lines <= 2 or chars < 16):
        return False

    # Keep at most one very large single-line box later as fallback.
    if lines <= 1 and chars < 24 and ratio < 0.08:
        return False

    return True


def pick_body_entries(shapes, canvas_area):
    raw_entries = []
    for shape in shapes:
        capacity = estimate_capacity(shape)
        if not capacity:
            continue
        raw_entries.append({
            "lines": capacity["lines"],
            "chars": capacity["chars"],
            "area": shape_area(shape),
            "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        })
    raw_entries.sort(key=lambda item: item["area"], reverse=True)

    multi_line = []
    single_line = []
    for entry in raw_entries:
        if not should_keep_body_entry(entry, canvas_area):
            continue
        if entry["lines"] >= 2:
            multi_line.append(entry)
        else:
            single_line.append(entry)

    selected = multi_line[:MAX_BODY_COLUMNS]
    if len(selected) < 2 and single_line:
        selected.append(single_line[0])

    for item in selected:
        item.pop("area", None)
        item.pop("is_placeholder", None)
    return selected


def analyze_shape_set(shapes, canvas_area):
    title_shapes = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title_shapes.append(shape)

    placeholders = collect_placeholder_shapes(shapes)
    titles, bodies, has_image = classify_placeholders(placeholders, title_shapes)

    if not bodies:
        # fallback to generic text shapes (exclude title shapes)
        bodies = collect_text_shapes(shapes, exclude=title_shapes)

    body_entries = pick_body_entries(bodies, canvas_area)

    title_capacity = None
    if title_shapes:
        title_capacity = estimate_capacity(title_shapes[0])
    title_chars = title_capacity["chars"] if title_capacity else 18

    return {
        "title_count": len(title_shapes),
        "body_columns": max(1, len(body_entries)) if body_entries else 1,
        "body_capacity": body_entries,
        "title_chars": title_chars,
        "has_image": has_image,
    }


def summarize(entries):
    if not entries:
        return {"lines": 4, "chars": 18}
    total_lines = 0
    total_chars = 0
    count = 0
    for entry in entries:
        total_lines += entry.get("lines", 0)
        total_chars += entry.get("chars", 0)
        count += 1
    if count == 0:
        return {"lines": 4, "chars": 18}
    avg_lines = max(2, int(total_lines / count))
    avg_chars = max(8, int(total_chars / count))
    return {"lines": avg_lines, "chars": avg_chars}


def build_summary(pres):
    canvas_area = int(pres.slide_width) * int(pres.slide_height)
    slides_summary = []
    all_body_caps = []
    for idx, slide in enumerate(pres.slides):
        info = analyze_shape_set(slide.shapes, canvas_area)
        for cap in info["body_capacity"]:
            all_body_caps.append(cap)
        slides_summary.append({
            "i": idx,
            "t": info["title_count"],
            "b": info["body_columns"],
            "cap": [{"lines": c["lines"], "chars": c["chars"]} for c in info["body_capacity"]],
            "titleChars": info["title_chars"],
            "img": 1 if info["has_image"] else 0,
        })

    layouts_summary = []
    for idx, layout in enumerate(pres.slide_layouts):
        info = analyze_shape_set(layout.shapes, canvas_area)
        layouts_summary.append({
            "i": idx,
            "t": info["title_count"],
            "b": info["body_columns"],
            "cap": [{"lines": c["lines"], "chars": c["chars"]} for c in info["body_capacity"]],
            "titleChars": info["title_chars"],
            "img": 1 if info["has_image"] else 0,
        })

    default_body = summarize(all_body_caps)
    return {
        "slides": slides_summary,
        "layouts": layouts_summary,
        "defaultBody": default_body,
        "slideCount": len(pres.slides),
        "layoutCount": len(pres.slide_layouts),
    }


def main():
    parser = argparse.ArgumentParser(description="Analyze PPTX template structure.")
    parser.add_argument("--template", required=True, help="Path to PPTX template")
    parser.add_argument("--output", required=True, help="Path to output JSON")
    args = parser.parse_args()

    template_path = Path(args.template)
    output_path = Path(args.output)

    if not template_path.exists():
        print(f"Template file not found: {template_path}", file=sys.stderr)
        return 2
    if not zipfile.is_zipfile(template_path):
        print(f"Template file is not a valid .pptx: {template_path}", file=sys.stderr)
        return 4

    try:
        pres = Presentation(str(template_path))
    except Exception as exc:
        print(f"Failed to load presentation: {exc}", file=sys.stderr)
        return 5

    stat = template_path.stat()
    analysis = {
        "version": ANALYSIS_VERSION,
        "template": {
            "path": str(template_path),
            "mtime": int(stat.st_mtime),
            "size": int(stat.st_size),
        },
        "summary": build_summary(pres),
    }

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", encoding="utf-8") as handle:
        json.dump(analysis, handle, ensure_ascii=False)
    return 0


if __name__ == "__main__":
    sys.exit(main())
