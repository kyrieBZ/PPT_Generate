#!/usr/bin/env python3
"""
analyze_style.py — PPT 风格迁移分析脚本 v2（方案二：截图 + 多模态 AI）

流程：
  1. LibreOffice headless 将 .pptx 转为 PDF
  2. pdftoppm 将 PDF 前 N 页转为 PNG 截图
  3. PIL quantize 从截图中提取主调色板（无需 sklearn）
  4. 将截图 base64 + 结构化 prompt 传给 Qwen-VL，
     让多模态 AI 以 JSON 输出精确的 StyleSpec
  5. 回退：若 Qwen-VL 不可用，仅返回 PIL 色彩聚类结果

同时保留 Theme XML 提取作为字体/颜色双重校验：
  - Theme XML 负责字体和标准色槽（精确无歧义）
  - Qwen-VL 负责语义化布局、视觉调性等难以用规则提取的信息

用法：
  python3 analyze_style.py --input ref.pptx [--api-key <key>] [--slides 3]

输出 JSON：
{
  "palette": { "primary": "#...", "secondary": "#...", "accent": "#...",
               "background": "#...", "text_dark": "#...", "text_light": "#..." },
  "typography": { "title_font": "...", "body_font": "...", "title_size": 36, "body_size": 18 },
  "layout": { "title_position": "top", "has_accent_bar": true,
               "has_fullscreen_bg": false, "geometric_style": "sharp", "content_ratio": 0.75 },
  "visual_style": "corporate_blue_minimal",
  "style_description": "深蓝主色调，简洁商务风格，标题居左，白色内容区",
  "theme_colors": { "dk1": "#...", ... },
  "preview_palette": ["#...", ...],
  "sample_slides": 3,
  "source": "vision_ai"
}
"""

import argparse
import base64
import io
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from collections import Counter

try:
    from PIL import Image
    from lxml import etree
    from pptx import Presentation
except ImportError as exc:
    print(json.dumps({"error": f"依赖缺失: pip install python-pptx lxml Pillow\n({exc})"}))
    sys.exit(1)

# ---------- 常量 -------------------------------------------------------------

_THEME_NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
_THEME_SLOTS = ['dk1', 'lt1', 'dk2', 'lt2',
                'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6']

# Qwen-VL API endpoint（DashScope OpenAI compatible）
_QWEN_VL_ENDPOINT = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
_QWEN_VL_MODEL    = "qwen-vl-plus"

# LibreOffice / pdftoppm
_SOFFICE = shutil.which("soffice") or shutil.which("libreoffice") or "libreoffice"
_PDFTOPPM = shutil.which("pdftoppm") or "pdftoppm"


# ---------- 颜色工具 ---------------------------------------------------------

def hex_to_rgb(h):
    h = h.lstrip('#')
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

def rgb_to_hex(r, g, b):
    return "#{:02X}{:02X}{:02X}".format(int(r), int(g), int(b))

def luminance(r, g, b):
    def lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)

def is_dark(r, g, b):
    return luminance(r, g, b) < 0.35

def hex_distance(h1, h2):
    r1, g1, b1 = hex_to_rgb(h1)
    r2, g2, b2 = hex_to_rgb(h2)
    return ((r1-r2)**2 + (g1-g2)**2 + (b1-b2)**2) ** 0.5

def is_neutral(h):
    r, g, b = hex_to_rgb(h)
    return max(r, g, b) - min(r, g, b) < 15

def deduplicate_colors(colors, min_dist=25):
    unique = []
    for c in colors:
        if all(hex_distance(c, u) >= min_dist for u in unique):
            unique.append(c)
    return unique


# ---------- Step 1: LibreOffice PPTX → PDF ----------------------------------

def pptx_to_pdf(pptx_path: str, out_dir: str) -> str | None:
    """将 pptx 转为 PDF，返回 PDF 路径；失败返回 None。"""
    try:
        result = subprocess.run(
            [_SOFFICE, "--headless", "--convert-to", "pdf",
             "--outdir", out_dir, pptx_path],
            capture_output=True, timeout=60
        )
        if result.returncode != 0:
            return None
        stem = Path(pptx_path).stem
        pdf_path = Path(out_dir) / (stem + ".pdf")
        return str(pdf_path) if pdf_path.exists() else None
    except Exception:
        return None


# ---------- Step 2: PDF → PNG 截图 ------------------------------------------

def pdf_to_slides(pdf_path: str, out_dir: str, max_pages: int = 3) -> list[str]:
    """
    用 pdftoppm 将 PDF 前 max_pages 页转为 PNG。
    返回 PNG 路径列表（按页序排序）。
    """
    prefix = str(Path(out_dir) / "slide")
    try:
        subprocess.run(
            [_PDFTOPPM, "-png", "-r", "72",
             "-f", "1", "-l", str(max_pages),
             pdf_path, prefix],
            capture_output=True, timeout=60
        )
    except Exception:
        pass
    pngs = sorted(Path(out_dir).glob("slide-*.png"))
    return [str(p) for p in pngs[:max_pages]]


# ---------- Step 3: PIL 色彩聚类 ---------------------------------------------

def extract_palette_from_images(img_paths: list[str], n_colors: int = 12) -> list[str]:
    """
    将多张截图合并后，用 PIL quantize 提取 n_colors 个主色，
    过滤近黑 / 近白后返回去重的 hex 颜色列表（按频次降序）。
    """
    if not img_paths:
        return []

    # 拼接所有截图为一张宽图（降采样以提速）
    frames = []
    for p in img_paths:
        try:
            img = Image.open(p).convert("RGB")
            img.thumbnail((480, 270))  # 降分辨率：保证速度
            frames.append(img)
        except Exception:
            pass

    if not frames:
        return []

    total_w = sum(f.width for f in frames)
    max_h   = max(f.height for f in frames)
    combined = Image.new("RGB", (total_w, max_h), (255, 255, 255))
    x = 0
    for f in frames:
        combined.paste(f, (x, 0))
        x += f.width

    # PIL quantize：快速八叉树量化
    quantized = combined.quantize(colors=n_colors, method=Image.Quantize.FASTOCTREE)
    pal_raw   = quantized.getpalette()[:n_colors * 3]

    # 统计各色出现频次
    data = quantized.getdata()
    freq = Counter(data)

    colors_by_freq = []
    for idx, count in freq.most_common(n_colors):
        r = pal_raw[idx * 3]
        g = pal_raw[idx * 3 + 1]
        b = pal_raw[idx * 3 + 2]
        colors_by_freq.append((rgb_to_hex(r, g, b), count))

    # 过滤极端中性色（纯白 / 近黑）和 near-duplicate
    filtered = [c for c, _ in colors_by_freq]
    return deduplicate_colors(filtered, min_dist=20)


def build_palette_from_colors(colors: list[str]) -> dict:
    """将原始颜色列表按语义映射为 palette 字典。"""
    dark   = [c for c in colors if is_dark(*hex_to_rgb(c))]
    light  = [c for c in colors if luminance(*hex_to_rgb(c)) > 0.75]
    mid    = [c for c in colors if not is_dark(*hex_to_rgb(c))
              and luminance(*hex_to_rgb(c)) <= 0.75
              and not is_neutral(c)]

    non_neutral_dark = [c for c in dark if not is_neutral(c)]
    primary    = non_neutral_dark[0] if non_neutral_dark else (dark[0] if dark else "#1E3A5F")
    background = light[0] if light else "#F8F9FA"
    secondary  = "#2D6A4F"
    if len(non_neutral_dark) > 1:
        secondary = non_neutral_dark[1]
    elif mid:
        secondary = mid[0]

    accent = "#F4A261"
    best_dist = 0
    for c in (mid or non_neutral_dark)[:8]:
        d = hex_distance(c, primary)
        if d > best_dist:
            best_dist = d
            accent = c

    text_dark  = dark[0] if dark else "#0D1B2A"
    text_light = light[0] if light else "#FFFFFF"

    return {
        "primary":    primary,
        "secondary":  secondary,
        "accent":     accent,
        "background": background,
        "text_dark":  text_dark,
        "text_light": text_light,
    }


# ---------- Step 4: Theme XML（字体 + 色槽校验）------------------------------

def extract_theme_xml(pptx_path: str) -> dict:
    """读取 ppt/theme/theme1.xml，提取 10 色槽和 major/minor 字体。"""
    result = {}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            theme_files = sorted(
                n for n in z.namelist()
                if n.startswith('ppt/theme/') and n.endswith('.xml')
            )
            if not theme_files:
                return {}
            root = etree.fromstring(z.read(theme_files[0]))

            # 色槽
            clr = root.find('.//a:clrScheme', _THEME_NS)
            if clr is not None:
                colors = {}
                for slot in _THEME_SLOTS:
                    child = clr.find(f'a:{slot}', _THEME_NS)
                    if child is not None:
                        srgb = child.find('a:srgbClr', _THEME_NS)
                        sys_clr = child.find('a:sysClr', _THEME_NS)
                        val = None
                        if srgb is not None:
                            val = srgb.get('val', '')
                        elif sys_clr is not None:
                            val = sys_clr.get('lastClr', '')
                        if val:
                            colors[slot] = '#' + val.upper()
                if colors:
                    result['theme_colors'] = colors

            # 字体
            fs = root.find('.//a:fontScheme', _THEME_NS)
            if fs is not None:
                maj = fs.find('a:majorFont/a:latin', _THEME_NS)
                min_ = fs.find('a:minorFont/a:latin', _THEME_NS)
                if maj is not None and not maj.get('typeface','').startswith('+'):
                    result['major_font'] = maj.get('typeface', '')
                if min_ is not None and not min_.get('typeface','').startswith('+'):
                    result['minor_font'] = min_.get('typeface', '')
    except Exception:
        pass
    return result


def scan_font_sizes(prs) -> tuple[int, int]:
    """从 Slide Master 和前 5 张幻灯片扫描字号。"""
    title_sizes, body_sizes = [], []
    for src in [prs.slide_master] + list(prs.slides)[:5]:
        for shape in src.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            try:
                ph = shape.placeholder_format
            except Exception:
                ph = None
            is_title = ph is not None and ph.idx in (0, 1)
            for para in shape.text_frame.paragraphs:
                try:
                    if para.font.size:
                        sz = round(para.font.size.pt)
                        (title_sizes if is_title else body_sizes).append(sz)
                except Exception:
                    pass
                for run in para.runs:
                    try:
                        if run.font.size:
                            sz = round(run.font.size.pt)
                            (title_sizes if is_title else body_sizes).append(sz)
                    except Exception:
                        pass

    def mc(lst, d):
        return Counter(lst).most_common(1)[0][0] if lst else d

    return mc(title_sizes, 36), mc(body_sizes, 18)


# ---------- Step 5: Qwen-VL 视觉风格分析 ------------------------------------

def _images_to_b64(img_paths: list[str]) -> list[str]:
    """将 PNG 路径转为 base64 data URI。"""
    result = []
    for p in img_paths:
        try:
            with open(p, 'rb') as f:
                b64 = base64.b64encode(f.read()).decode()
            result.append("data:image/png;base64," + b64)
        except Exception:
            pass
    return result


def analyze_style_with_vision(
    img_paths: list[str],
    palette_hint: dict,
    typography_hint: dict,
    api_key: str,
) -> dict | None:
    """
    将截图 + PIL 提取的调色板作为 hint 传给 Qwen-VL，
    让 AI 输出精确的 StyleSpec JSON。
    返回解析后的字典，失败返回 None。
    """
    try:
        import urllib.request
        import urllib.error
    except ImportError:
        return None

    images_b64 = _images_to_b64(img_paths)
    if not images_b64:
        return None

    # 构造 content 数组：图片 + 文字
    content = []
    for b64 in images_b64:
        content.append({"type": "image_url", "image_url": {"url": b64}})

    palette_str = json.dumps(palette_hint, ensure_ascii=False)
    typo_str    = json.dumps(typography_hint, ensure_ascii=False)

    prompt = f"""你是一位精通 PowerPoint 设计的视觉设计师。
以上是某份参考 PPT 的幻灯片截图（1~3 张）。
我已通过 PIL 色彩聚类初步提取了调色板，请你在此基础上进行精准分析和校正。

【PIL 初步提取结果（供参考，可修正）】
调色板: {palette_str}
字体:   {typo_str}

请仔细观察截图，输出以下 JSON，不要有任何额外文字或 markdown：
{{
  "palette": {{
    "primary":    "最主要的深色（标题栏/色块背景色），必须是6位十六进制",
    "secondary":  "辅助色（副标题栏/装饰条），6位十六进制",
    "accent":     "强调色（高亮图标/按钮/装饰点），6位十六进制",
    "background": "幻灯片内容区背景色（通常较浅），6位十六进制",
    "text_dark":  "浅色背景上的文字颜色，6位十六进制",
    "text_light": "深色背景上的文字颜色，6位十六进制"
  }},
  "typography": {{
    "title_font": "标题字体名，只能是 PowerPoint 内置字体如 Calibri/Arial/Georgia 等",
    "body_font":  "正文字体名",
    "title_size": 标题字号数字,
    "body_size":  正文字号数字
  }},
  "layout": {{
    "title_position":    "top 或 center",
    "has_accent_bar":    是否有横/纵装饰色条（true/false）,
    "has_fullscreen_bg": 是否有全屏图片背景（true/false）,
    "geometric_style":   "sharp（直角）或 rounded（圆角）",
    "content_ratio":     内容区域占页面的比例（0~1之间的小数）
  }},
  "visual_style": "用3~5个英文单词描述整体风格，如 corporate_blue_minimal",
  "style_description": "用一句中文描述这份PPT的视觉风格特点，30字以内"
}}

注意：
- palette 中每个值都必须是 # 开头的6位十六进制颜色
- primary 不能是纯白 #FFFFFF
- 若截图信息不足，请基于 PIL 提取结果合理推断"""

    content.append({"type": "text", "text": prompt})

    body = {
        "model": _QWEN_VL_MODEL,
        "messages": [
            {"role": "user", "content": content}
        ],
        "max_tokens": 800,
    }

    try:
        payload = json.dumps(body).encode('utf-8')
        req = urllib.request.Request(
            _QWEN_VL_ENDPOINT,
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=60) as resp:
            resp_data = json.loads(resp.read().decode('utf-8'))
        text = resp_data["choices"][0]["message"]["content"]
    except Exception as e:
        return None

    # 解析 JSON
    start = text.find('{')
    end   = text.rfind('}')
    if start == -1 or end == -1:
        return None
    try:
        return json.loads(text[start:end+1])
    except Exception:
        return None


# ---------- 布局分析（规则回退，当 VL 失败时使用）---------------------------

def analyze_layout_rules(slides, slide_w, slide_h) -> dict:
    title_pos_counts = Counter()
    has_bar = False
    has_fullscreen_bg = False
    rounded_count = 0
    sharp_count = 0
    content_ratios = []

    for slide in slides:
        content_shapes = 0
        total_shapes = 0

        try:
            bg = slide.background.fill
            if bg.type is not None and str(bg.type) in ('BACKGROUND', 'PICTURE', '3'):
                has_fullscreen_bg = True
        except Exception:
            pass

        for shape in slide.shapes:
            total_shapes += 1
            try:
                ph = shape.placeholder_format
            except Exception:
                ph = None

            if ph is not None and ph.idx in (0, 1):
                try:
                    frac = shape.top / slide_h
                    title_pos_counts['top' if frac < 0.3 else 'center'] += 1
                except Exception:
                    title_pos_counts['top'] += 1
            else:
                content_shapes += 1

            if not has_bar:
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        w_frac = shape.width / slide_w
                        h_frac = shape.height / slide_h
                        if (w_frac > 0.5 and h_frac < 0.1) or (h_frac > 0.5 and w_frac < 0.1):
                            has_bar = True
                except Exception:
                    pass

            try:
                from pptx.oxml.ns import qn as _qn
                prst = shape.element.find('.//' + _qn('a:prstGeom'))
                if prst is not None:
                    name = prst.get('prst', '')
                    if 'round' in name.lower():
                        rounded_count += 1
                    else:
                        sharp_count += 1
            except Exception:
                pass

        content_ratios.append(content_shapes / total_shapes if total_shapes else 0.7)

    title_pos = title_pos_counts.most_common(1)[0][0] if title_pos_counts else 'top'
    avg_ratio = round(sum(content_ratios) / len(content_ratios), 2) if content_ratios else 0.7
    geometric_style = 'rounded' if rounded_count > sharp_count else 'sharp'

    return {
        "title_position":    title_pos,
        "has_accent_bar":    has_bar,
        "has_fullscreen_bg": has_fullscreen_bg,
        "geometric_style":   geometric_style,
        "content_ratio":     avg_ratio,
    }


# ---------- 主分析入口 -------------------------------------------------------

def analyze(pptx_path: str, api_key: str = "", max_slides: int = 3) -> dict:
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"error": f"无法读取 PPTX 文件: {e}"}

    slides = list(prs.slides)
    if not slides:
        return {"error": "PPT 文件中没有幻灯片"}

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    sample_slides = slides[:min(len(slides), max_slides)]

    # ---- Theme XML：提取精确字体 + 色槽（不受 VL 影响，始终执行）----
    theme_info   = extract_theme_xml(pptx_path)
    theme_colors = theme_info.get('theme_colors', {})
    title_size, body_size = scan_font_sizes(prs)

    # 字体：Theme fontScheme > txStyles > 默认
    title_font = theme_info.get('major_font', 'Calibri') or 'Calibri'
    body_font  = theme_info.get('minor_font', 'Calibri') or 'Calibri'

    typography_hint = {
        "title_font": title_font,
        "body_font":  body_font,
        "title_size": title_size,
        "body_size":  body_size,
    }

    # ---- 截图流程 ----
    img_paths = []
    with tempfile.TemporaryDirectory() as tmp_dir:
        pdf_path = pptx_to_pdf(pptx_path, tmp_dir)
        if pdf_path:
            img_paths = pdf_to_slides(pdf_path, tmp_dir, max_pages=max_slides)

        # ---- PIL 色彩聚类（始终执行，作为 VL 的 hint 或 fallback）----
        raw_colors   = extract_palette_from_images(img_paths) if img_paths else []
        palette_hint = build_palette_from_colors(raw_colors) if raw_colors else {}

        # 若 Theme XML 有标准色槽，用它校正 palette_hint（dk2→primary 等）
        if theme_colors:
            def tc(slot, fallback):
                return theme_colors.get(slot, fallback)
            # 用主题色覆盖 PIL 的粗糙结果（主题色最权威）
            palette_hint["primary"]    = tc('dk2', palette_hint.get('primary', '#1E3A5F'))
            palette_hint["background"] = tc('lt2', palette_hint.get('background', '#F8F9FA'))
            r, g, b = hex_to_rgb(palette_hint["background"])
            if luminance(r, g, b) < 0.5:
                palette_hint["background"] = tc('lt1', '#FFFFFF')
            palette_hint["text_dark"]  = tc('dk1', palette_hint.get('text_dark', '#000000'))
            palette_hint["text_light"] = tc('lt1', palette_hint.get('text_light', '#FFFFFF'))
            palette_hint["secondary"]  = tc('accent1', palette_hint.get('secondary', '#2D6A4F'))
            # accent: 从 accent1~6 里找与 primary 差距最大的
            best, best_d = palette_hint.get('accent', '#F4A261'), 0
            for slot in ['accent1','accent2','accent3','accent4','accent5','accent6']:
                c = theme_colors.get(slot)
                if c:
                    d = hex_distance(c, palette_hint['primary'])
                    if d > best_d:
                        best_d, best = d, c
            palette_hint["accent"] = best

        # ---- Qwen-VL 视觉分析 ----
        vl_result = None
        if api_key and img_paths:
            vl_result = analyze_style_with_vision(
                img_paths, palette_hint, typography_hint, api_key
            )

    # ---- 合并结果 ----
    if vl_result:
        # VL 成功：以 VL 结果为主，Theme XML 字体为校验
        palette = vl_result.get('palette', palette_hint)
        # 确保所有颜色值有效
        for k, v in palette_hint.items():
            if k not in palette or not isinstance(palette[k], str) or not palette[k].startswith('#'):
                palette[k] = v

        typography = vl_result.get('typography', typography_hint)
        # 字体用 Theme XML 结果覆盖（VL 在字体上不如 XML 精确）
        if title_font and title_font != 'Calibri':
            typography['title_font'] = title_font
        if body_font and body_font != 'Calibri':
            typography['body_font'] = body_font
        typography['title_size'] = typography.get('title_size') or title_size
        typography['body_size']  = typography.get('body_size')  or body_size

        layout = vl_result.get('layout', {})
        visual_style     = vl_result.get('visual_style', '')
        style_description = vl_result.get('style_description', '')
        source = "vision_ai"
    else:
        # Fallback：PIL 聚类 + Theme XML 规则
        palette = palette_hint if palette_hint else {
            "primary": "#1E3A5F", "secondary": "#2D6A4F", "accent": "#F4A261",
            "background": "#F8F9FA", "text_dark": "#0D1B2A", "text_light": "#FFFFFF"
        }
        typography = typography_hint
        layout = analyze_layout_rules(sample_slides, slide_w, slide_h)
        visual_style = ""
        style_description = ""
        source = "theme_xml_pil" if theme_colors else "frequency_fallback"

    preview_palette = list(dict.fromkeys([
        palette["primary"], palette["secondary"], palette["accent"],
        palette["background"], palette["text_dark"], palette["text_light"],
    ]))

    result = {
        "palette":           palette,
        "typography":        typography,
        "layout":            layout,
        "preview_palette":   preview_palette,
        "sample_slides":     len(sample_slides),
        "source":            source,
    }
    if visual_style:
        result["visual_style"] = visual_style
    if style_description:
        result["style_description"] = style_description
    if theme_colors:
        result["theme_colors"] = theme_colors

    return result


# ---------- CLI --------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Analyze PPTX style (Vision AI + Theme XML)")
    parser.add_argument("--input",    required=True, help="Path to reference .pptx")
    parser.add_argument("--api-key",  default="",    help="Qwen API key for Vision AI")
    parser.add_argument("--slides",   type=int, default=3, help="Number of slides to capture (default 3)")
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(json.dumps({"error": f"文件不存在: {args.input}"}))
        sys.exit(1)

    result = analyze(args.input, api_key=args.api_key, max_slides=args.slides)
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
