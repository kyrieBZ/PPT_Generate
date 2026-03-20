#!/usr/bin/env python3
import argparse
import json
import re
import subprocess
import sys
from pathlib import Path
import zipfile
import os
import datetime
import tempfile
import urllib.request

try:
    from pptx import Presentation
    from pptx.enum.dml import MSO_COLOR_TYPE
    from pptx.enum.shapes import PP_PLACEHOLDER
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
    except ImportError:
        MSO_AUTO_SIZE = None
    try:
        from pptx.util import Inches
    except ImportError:
        Inches = None
except ImportError as exc:
    print("python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    raise


EMU_PER_PT = 12700.0
MIN_FONT_PT = 12.0
FONT_STEP_PT = 2.0
MAX_CONTINUATION_SLIDES = 10

_TEMPLATE_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_PRESENTATION_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def open_presentation(template_path):
    """Open a PPTX or POTX file as a python-pptx Presentation.

    Files saved as PowerPoint templates (.potx) have a different content-type
    entry in [Content_Types].xml that python-pptx refuses to load.  This
    function detects that case, rewrites the content-type in memory, and
    returns the loaded Presentation so callers never see the error.
    """
    if not zipfile.is_zipfile(template_path):
        return Presentation(template_path)

    # Check if the content-type needs patching (read [Content_Types].xml from zip)
    with zipfile.ZipFile(template_path, "r") as _zcheck:
        try:
            _ct_data = _zcheck.read("[Content_Types].xml")
        except KeyError:
            _ct_data = b""

    if _TEMPLATE_CONTENT_TYPE.encode("utf-8") not in _ct_data:
        return Presentation(template_path)

    # Patch the content-type in-place within the zip so python-pptx accepts it
    buf = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    buf.close()
    try:
        with zipfile.ZipFile(template_path, "r") as zin:
            with zipfile.ZipFile(buf.name, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = data.replace(
                            _TEMPLATE_CONTENT_TYPE.encode("utf-8"),
                            _PRESENTATION_CONTENT_TYPE.encode("utf-8"),
                        )
                    zout.writestr(item, data)
        pres = Presentation(buf.name)
    finally:
        try:
            os.remove(buf.name)
        except Exception:
            pass
    return pres


def shape_token(shape):
    if shape is None:
        return None
    shape_id = getattr(shape, "shape_id", None)
    if shape_id is not None:
        return ("shape_id", shape_id)
    element = getattr(shape, "element", None)
    if element is not None:
        return ("element_id", id(element))
    return ("object_id", id(shape))


def same_shape(left, right):
    return shape_token(left) == shape_token(right)


def delete_slide(pres, index):
    slide_id_list = pres.slides._sldIdLst  # pylint: disable=protected-access
    slide_ids = list(slide_id_list)
    if index < 0 or index >= len(slide_ids):
        return
    slide_id_list.remove(slide_ids[index])


def find_title_placeholder(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return shape
    if slide.shapes.title:
        return slide.shapes.title
    return None


def collect_body_placeholders(slide, title_shape):
    body_types = []
    for name in ("BODY", "CONTENT", "TEXT", "OBJECT", "SUBTITLE"):
        if hasattr(PP_PLACEHOLDER, name):
            body_types.append(getattr(PP_PLACEHOLDER, name))
    body_types = tuple(body_types)
    placeholders = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if title_shape is not None and same_shape(shape, title_shape):
            continue
        placeholder_type = shape.placeholder_format.type
        if body_types and placeholder_type in body_types:
            placeholders.append(shape)
    return placeholders


def collect_non_title_text_placeholders(slide, title_shape):
    placeholders = []
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if title_shape is not None and same_shape(shape, title_shape):
            continue
        if getattr(shape, "has_text_frame", False):
            placeholders.append(shape)
    return placeholders


def collect_text_shapes(slide, title_shape):
    candidates = []
    small = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if title_shape is not None and same_shape(shape, title_shape):
            continue
        capacity = estimate_capacity(shape)
        if capacity and capacity.get("lines", 0) >= 2 and capacity.get("chars", 0) >= 8:
            candidates.append(shape)
        else:
            small.append(shape)
    return candidates if candidates else small


def collect_placeholder_shapes(shapes):
    placeholders = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        placeholders.append(shape)
    return placeholders


def classify_placeholders(placeholders, title_shapes):
    body_types = []
    for name in ("BODY", "CONTENT", "TEXT", "OBJECT", "SUBTITLE"):
        if hasattr(PP_PLACEHOLDER, name):
            body_types.append(getattr(PP_PLACEHOLDER, name))
    body_types = tuple(body_types)

    image_types = []
    for name in ("PICTURE", "CHART", "TABLE", "CLIP_ART", "MEDIA"):
        if hasattr(PP_PLACEHOLDER, name):
            image_types.append(getattr(PP_PLACEHOLDER, name))
    image_types = tuple(image_types)

    bodies = []
    images = []
    has_image = False
    for shape in placeholders:
        if any(same_shape(shape, item) for item in title_shapes):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in body_types:
            bodies.append(shape)
        elif ph_type in image_types:
            images.append(shape)
            has_image = True
        else:
            if getattr(shape, "has_text_frame", False):
                bodies.append(shape)
    return bodies, images, has_image


def collect_image_placeholders(slide, title_shape):
    placeholders = collect_placeholder_shapes(slide.shapes)
    _, images, _ = classify_placeholders(placeholders, [title_shape] if title_shape else [])
    return images


def analyze_shapes(shapes):
    title_shapes = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title_shapes.append(shape)

    placeholders = collect_placeholder_shapes(shapes)
    bodies, _, has_image = classify_placeholders(placeholders, title_shapes)
    if not bodies:
        bodies = []
        for shape in shapes:
            if any(same_shape(shape, item) for item in title_shapes):
                continue
            if getattr(shape, "has_text_frame", False):
                bodies.append(shape)

    body_caps = []
    total_lines = 0
    for shape in bodies:
        cap = estimate_capacity(shape)
        if cap:
            body_caps.append(cap)
            total_lines += cap.get("lines", 0)

    return {
        "title_count": len(title_shapes),
        "body_columns": max(1, len(body_caps)) if body_caps else 1,
        "has_image": has_image,
        "total_lines": total_lines,
    }


def parse_image_prompts(slide_data):
    prompts = []
    for key in ("imagePrompts", "image_prompts"):
        value = slide_data.get(key, [])
        if isinstance(value, list):
            prompts.extend([str(item).strip() for item in value if str(item).strip()])
    return prompts


def download_image(url, dest_path):
    with urllib.request.urlopen(url, timeout=12) as response:
        if response.status and response.status >= 400:
            raise RuntimeError(f"image download failed: {response.status}")
        data = response.read()
    with open(dest_path, "wb") as handle:
        handle.write(data)


def resolve_image_sources(slide_data, temp_files):
    sources = []
    for key in ("imagePaths", "image_paths"):
        value = slide_data.get(key, [])
        if isinstance(value, list):
            for item in value:
                if not item:
                    continue
                path = str(item)
                if os.path.exists(path):
                    sources.append(path)
    for key in ("imageUrls", "image_urls"):
        value = slide_data.get(key, [])
        if isinstance(value, list):
            for item in value:
                url = str(item).strip()
                if not url:
                    continue
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.close()
                try:
                    download_image(url, tmp.name)
                    temp_files.append(tmp.name)
                    sources.append(tmp.name)
                except Exception as exc:
                    _log(f"image download failed: {exc}")
                    try:
                        os.remove(tmp.name)
                    except Exception:
                        pass
    return sources


def pick_layout_index(slide_data, layout_profiles, default_index):
    if not layout_profiles:
        return default_index

    bullet_groups = normalize_bullet_groups(
        slide_data.get("bulletGroups") or slide_data.get("bullet_groups") or []
    )
    bullets = slide_data.get("bullets", []) or []
    if not isinstance(bullets, list):
        bullets = [str(bullets)]
    bullets = [str(item).strip() for item in bullets if str(item).strip()]

    desired_columns = len(bullet_groups) if bullet_groups else 1
    if bullet_groups:
        desired_total_lines = sum(len(group) for group in bullet_groups)
    else:
        desired_total_lines = len(bullets)
    desired_has_image = bool(parse_image_prompts(slide_data))

    best_index = default_index
    best_score = None
    for profile in layout_profiles:
        score = 0
        score += abs(profile["body_columns"] - desired_columns) * 10
        if desired_has_image and not profile["has_image"]:
            score += 15
        elif not desired_has_image and profile["has_image"]:
            score += 3
        if profile["title_count"] == 0:
            score += 5
        if desired_total_lines and profile["total_lines"] < desired_total_lines:
            score += (desired_total_lines - profile["total_lines"]) * 2
        if profile["body_columns"] > desired_columns:
            score += (profile["body_columns"] - desired_columns) * 2

        if best_score is None or score < best_score:
            best_score = score
            best_index = profile["index"]
    return best_index


def pick_layout_index_from_guide(guide_page, layout_profiles, default_index, desired_has_image=False):
    """Choose layout index from layout_guide page (groups count, etc.). Used when payload has layout_guide."""
    if not layout_profiles:
        return default_index
    groups = guide_page.get("groups") if isinstance(guide_page, dict) else []
    desired_columns = len(groups) if groups else 1
    desired_total_lines = sum(
        int(g.get("max_bullets", 1)) for g in groups if isinstance(g, dict)
    ) if groups else 1

    best_index = default_index
    best_score = None
    for profile in layout_profiles:
        score = 0
        score += abs(profile["body_columns"] - desired_columns) * 10
        if desired_has_image and not profile["has_image"]:
            score += 15
        elif not desired_has_image and profile["has_image"]:
            score += 3
        if profile["title_count"] == 0:
            score += 5
        if desired_total_lines and profile["total_lines"] < desired_total_lines:
            score += (desired_total_lines - profile["total_lines"]) * 2
        if profile["body_columns"] > desired_columns:
            score += (profile["body_columns"] - desired_columns) * 2

        if best_score is None or score < best_score:
            best_score = score
            best_index = profile["index"]
    return best_index


def shape_position(shape):
    try:
        left = int(getattr(shape, "left", 0))
        top = int(getattr(shape, "top", 0))
    except Exception:
        left = 0
        top = 0
    return (left, top, -shape_area(shape))


def shape_area(shape):
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def emu_to_pt(value):
    try:
        return float(value) / EMU_PER_PT
    except Exception:
        return 0.0


def pick_font_pt(shape, default_pt=18.0):
    if shape is None or not shape.has_text_frame:
        return default_pt
    try:
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return default_pt
        paragraph = text_frame.paragraphs[0]
        if paragraph.runs:
            size = paragraph.runs[0].font.size
            if size is not None:
                return size.pt
        if paragraph.font and paragraph.font.size is not None:
            return paragraph.font.size.pt
    except Exception:
        return default_pt
    return default_pt


def estimate_capacity(shape):
    if shape is None or not shape.has_text_frame:
        return None
    font_pt = pick_font_pt(shape)
    return estimate_capacity_for_font(shape, font_pt)


def estimate_capacity_for_font(shape, font_pt):
    if shape is None or not shape.has_text_frame:
        return None
    width_pt = emu_to_pt(getattr(shape, "width", 0))
    height_pt = emu_to_pt(getattr(shape, "height", 0))
    text_frame = shape.text_frame
    margin_left = emu_to_pt(getattr(text_frame, "margin_left", 0))
    margin_right = emu_to_pt(getattr(text_frame, "margin_right", 0))
    margin_top = emu_to_pt(getattr(text_frame, "margin_top", 0))
    margin_bottom = emu_to_pt(getattr(text_frame, "margin_bottom", 0))
    usable_width = max(width_pt - margin_left - margin_right, font_pt)
    usable_height = max(height_pt - margin_top - margin_bottom, font_pt * 1.2)
    line_height = max(font_pt * 1.2, 10.0)
    max_lines = max(1, int(usable_height / line_height))
    char_width = max(font_pt * 0.9, 8.0)
    max_chars = max(4, int(usable_width / char_width))
    return {"lines": max_lines, "chars": max_chars}


def truncate_text(text, max_chars):
    if max_chars <= 0:
        return ""
    if text is None:
        return ""
    return str(text).strip()


def wrap_text(text, max_chars):
    if max_chars <= 0:
        return []
    if text is None:
        return []
    value = str(text).strip()
    if not value:
        return []
    try:
        import unicodedata
    except Exception:
        unicodedata = None

    def char_width(ch):
        if unicodedata is None:
            return 1
        if unicodedata.east_asian_width(ch) in ("W", "F"):
            return 2
        return 1

    break_chars = set(" ，。、；;,.!?！？:：")
    lines = []
    line_start = 0
    last_break = -1
    width = 0
    idx = 0
    length = len(value)
    while idx < length:
        ch = value[idx]
        w = char_width(ch)
        if ch in break_chars:
            last_break = idx
        if width + w > max_chars:
            if last_break >= line_start:
                line = value[line_start:last_break + 1].strip()
                line_start = last_break + 1
                idx = line_start
            else:
                line = value[line_start:idx].strip()
                line_start = idx
            if line:
                lines.append(line)
            last_break = -1
            width = 0
            continue
        width += w
        idx += 1
    if line_start < length:
        tail = value[line_start:].strip()
        if tail:
            lines.append(tail)
    return lines


def normalize_text_item(text):
    if text is None:
        return ""
    value = str(text).replace("\r", " ").replace("\n", " ").replace("\t", " ").strip()
    # Collapse repeated spaces.
    while "  " in value:
        value = value.replace("  ", " ")
    return value


def fit_bullets_to_capacity(bullets, capacity):
    if capacity is None:
        return bullets, []
    max_lines = max(1, capacity.get("lines", 1))
    max_chars = max(4, capacity.get("chars", 12))
    fitted = []
    overflow = []
    used_lines = 0
    for bullet in bullets:
        bullet = normalize_text_item(bullet)
        if not bullet:
            continue
        wrapped = wrap_text(bullet, max_chars)
        needed = max(1, len(wrapped))
        if used_lines + needed <= max_lines:
            fitted.append(bullet)
            used_lines += needed
        else:
            overflow.append(bullet)
    return fitted, overflow


def normalize_bullet_groups(raw_groups):
    groups = []
    if not isinstance(raw_groups, list):
        return groups
    for group in raw_groups:
        if not isinstance(group, list):
            continue
        items = [str(item).strip() for item in group if str(item).strip()]
        if items:
            groups.append(items)
    return groups
def copy_font(src, dest):
    dest.bold = src.bold
    dest.italic = src.italic
    dest.underline = src.underline
    dest.size = src.size
    dest.name = src.name
    try:
        color_type = getattr(src.color, "type", None)
        if color_type == MSO_COLOR_TYPE.RGB:
            dest.color.rgb = src.color.rgb
        elif hasattr(src.color, "theme_color") and src.color.theme_color is not None:
            dest.color.theme_color = src.color.theme_color
        if hasattr(src.color, "brightness") and src.color.brightness is not None:
            dest.color.brightness = src.color.brightness
    except Exception:
        pass


def prune_runs(paragraph, keep=1):
    runs = list(paragraph.runs)
    for run in runs[keep:]:
        run._r.getparent().remove(run._r)


def prune_paragraphs(text_frame, keep=1):
    paragraphs = list(text_frame.paragraphs)
    for para in paragraphs[keep:]:
        para._p.getparent().remove(para._p)


def clear_text(shape):
    if shape is None or not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = ""
        prune_runs(paragraph, keep=1)
    else:
        paragraph.text = ""
    prune_paragraphs(text_frame, keep=1)


def _log(message):
    if os.environ.get("PPT_BUILDER_LOG", "").lower() not in ("1", "true", "yes", "on"):
        return
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    print(f"[ppt_builder] {timestamp} {message}", file=sys.stderr)


def apply_text(shape, lines):
    return apply_text_with_font(shape, lines, None)


def apply_text_with_font(shape, lines, font_pt):
    if shape is None or not shape.has_text_frame:
        return
    if not lines:
        clear_text(shape)
        return
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    if MSO_AUTO_SIZE is not None:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if not text_frame.paragraphs:
        text_frame.add_paragraph()
    first_paragraph = text_frame.paragraphs[0]
    if first_paragraph.runs:
        first_run = first_paragraph.runs[0]
    else:
        first_run = first_paragraph.add_run()

    # Always clear existing placeholder text to avoid overlay.
    first_run.text = ""
    prune_runs(first_paragraph, keep=1)
    prune_paragraphs(text_frame, keep=1)

    if font_pt is not None:
        try:
            from pptx.util import Pt  # lazy import
            first_run.font.size = Pt(float(font_pt))
        except Exception:
            pass

    first_run.text = normalize_text_item(lines[0])

    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.level = first_paragraph.level
        run = paragraph.add_run()
        copy_font(first_run.font, run.font)
        run.text = normalize_text_item(line)


def filter_body_shapes(shapes):
    # Prefer shapes with meaningful capacity to avoid filling decorative boxes.
    candidates = []
    fallback = []
    for shape in shapes:
        cap = estimate_capacity(shape)
        if cap:
            fallback.append(shape)
            if cap.get("lines", 0) >= 2 and cap.get("chars", 0) >= 8:
                candidates.append(shape)
    return candidates if candidates else fallback


def fit_with_font_steps(shape, bullets):
    # Try reducing font size to fit more content; fall back to best-effort if still overflow.
    bullets = [normalize_text_item(b) for b in bullets if normalize_text_item(b)]
    if not bullets:
        return [], [], None
    base_font = pick_font_pt(shape)
    if base_font is None:
        base_font = 18.0
    tried = []
    font = float(base_font)
    while font >= MIN_FONT_PT - 0.1:
        cap = estimate_capacity_for_font(shape, font)
        fitted, overflow = fit_bullets_to_capacity(bullets, cap)
        tried.append((len(overflow), fitted, overflow, font))
        if not overflow:
            return fitted, overflow, font
        font -= FONT_STEP_PT
    tried.sort(key=lambda t: t[0])
    best = tried[0]
    return best[1], best[2], best[3]


def cleanup_unused_text(slide, title_shape, used_placeholder_tokens, clean_template_text):
    def is_system_placeholder(shape):
        if not getattr(shape, "is_placeholder", False):
            return False
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            return False
        return ph_type in (
            PP_PLACEHOLDER.SLIDE_NUMBER,
            PP_PLACEHOLDER.FOOTER,
            PP_PLACEHOLDER.DATE,
            PP_PLACEHOLDER.HEADER,
        )

    cleared_placeholders = 0
    cleared_textboxes = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if title_shape is not None and same_shape(shape, title_shape):
            continue
        if shape_token(shape) in used_placeholder_tokens:
            continue
        if getattr(shape, "is_placeholder", False):
            if is_system_placeholder(shape):
                continue
            clear_text(shape)
            cleared_placeholders += 1
        elif clean_template_text:
            text = ""
            try:
                text = shape.text_frame.text
            except Exception:
                text = ""
            if text.strip():
                clear_text(shape)
                cleared_textboxes += 1
    return cleared_placeholders, cleared_textboxes


def fill_slide_text(slide, title, bullets, bullet_groups):
    title_shape = find_title_placeholder(slide)
    if title_shape is not None:
        apply_text(title_shape, [title] if title else [])

    placeholders = collect_body_placeholders(slide, title_shape)
    if not placeholders:
        placeholders = collect_non_title_text_placeholders(slide, title_shape)
    if not placeholders:
        placeholders = collect_text_shapes(slide, title_shape)

    placeholders = filter_body_shapes(placeholders)
    if placeholders:
        placeholders.sort(key=shape_position)

    used_placeholder_tokens = set()
    overflow = []

    bullets = [normalize_text_item(b) for b in (bullets or []) if normalize_text_item(b)]

    if placeholders and bullet_groups:
        groups = [[normalize_text_item(b) for b in group if normalize_text_item(b)] for group in bullet_groups]
        groups = [g for g in groups if g]
        # If model provides more groups than available placeholders, merge the tail.
        if len(groups) > len(placeholders) and placeholders:
            merged = []
            for group in groups[:len(placeholders) - 1]:
                merged.append(group)
            tail = []
            for group in groups[len(placeholders) - 1:]:
                tail.extend(group)
            merged.append(tail)
            groups = merged

        carry = []
        for i, shape in enumerate(placeholders):
            used_placeholder_tokens.add(shape_token(shape))
            incoming = []
            if i < len(groups):
                incoming = groups[i]
            elif carry:
                incoming = carry
                carry = []
            if not incoming:
                clear_text(shape)
                continue
            fitted, extra, font_pt = fit_with_font_steps(shape, incoming)
            if extra:
                carry.extend(extra)
            apply_text_with_font(shape, fitted, font_pt)
        overflow.extend(carry)
        if len(groups) > len(placeholders):
            for group in groups[len(placeholders):]:
                overflow.extend(group)
    elif placeholders:
        remaining = bullets
        for shape in placeholders:
            used_placeholder_tokens.add(shape_token(shape))
            if not remaining:
                clear_text(shape)
                continue
            fitted, remaining, font_pt = fit_with_font_steps(shape, remaining)
            apply_text_with_font(shape, fitted, font_pt)
        overflow = remaining
    else:
        overflow = bullets

    return overflow, title_shape, used_placeholder_tokens


def expand_bullets(bullets):
    if len(bullets) != 1:
        return bullets
    text = bullets[0]
    separators = ["。", "；", ";", "、"]
    for sep in separators:
        if sep in text:
            parts = [part.strip() for part in text.split(sep) if part.strip()]
            if len(parts) > 1:
                return parts
    return bullets


def split_bullets(bullets, buckets):
    if buckets <= 1:
        return [bullets]
    bullets = expand_bullets(bullets)
    if len(bullets) <= buckets:
        result = []
        for idx in range(buckets):
            if idx < len(bullets):
                result.append([bullets[idx]])
            else:
                result.append([])
        return result
    base = len(bullets) // buckets
    extra = len(bullets) % buckets
    result = []
    cursor = 0
    for idx in range(buckets):
        size = base + (1 if idx < extra else 0)
        result.append(bullets[cursor:cursor + size])
        cursor += size
    return result


# Placeholder-like patterns to detect unfilled template text (3.4 post-build QA)
_PLACEHOLDER_PATTERNS = re.compile(
    r"xxxx|lorem\s+ipsum|Click to add|此页版式|点击添加|\[.*占位.*\]",
    re.IGNORECASE,
)


def run_placeholder_qa(output_path, strict_qa=False):
    """Run markitdown on output PPTX and grep for placeholder-like text; log warnings or exit non-zero if --strict-qa."""
    if not output_path or not os.path.isfile(output_path):
        return
    try:
        out = subprocess.run(
            [sys.executable, "-m", "markitdown", output_path],
            capture_output=True,
            text=True,
            timeout=60,
        )
    except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
        print(f"PPT_QA: skipped (markitdown not run): {e}", file=sys.stderr)
        return
    text = (out.stdout or "") + (out.stderr or "")
    hits = _PLACEHOLDER_PATTERNS.findall(text)
    if hits:
        unique = list(dict.fromkeys(hits))
        msg = f"PPT_QA: possible placeholder left in output: {unique}"
        print(msg, file=sys.stderr)
        if strict_qa:
            sys.exit(1)
    return


def build_presentation(template_path, output_path, payload, strict_qa=False):
    _log(f"build_presentation template={template_path} output={output_path}")
    strict_qa = strict_qa or bool(payload.get("strictQa", False))
    pres = open_presentation(template_path)

    slides = payload.get("slides", [])
    layout_mode = payload.get("layoutMode", "template")
    layout_guide = payload.get("layout_guide")  # optional: list of per-slide layout constraints from backend
    if not isinstance(layout_guide, list):
        layout_guide = []
    clean_template_text = bool(payload.get("cleanTemplateText", False))
    layouts = pres.slide_layouts
    temp_files = []
    layout_profiles = []
    fallback_profiles = []
    for idx, layout in enumerate(layouts):
        profile = analyze_shapes(layout.shapes)
        profile["index"] = idx
        profile["placeholder_count"] = len(layout.placeholders)
        fallback_profiles.append(profile)
        # Layouts without placeholders often create blank slides after add_slide.
        if profile["placeholder_count"] > 0:
            layout_profiles.append(profile)
    if not layout_profiles:
        layout_profiles = fallback_profiles

    existing_count = len(pres.slides)
    # In template mode with existing slides: fill them to preserve background; only add new if needed.
    use_existing_slides = layout_mode == "template" and existing_count > 0

    for idx, slide_data in enumerate(slides):
        _log(f"slide[{idx}] begin")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", []) or []
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        bullets = [str(item) for item in bullets if str(item).strip()]
        bullet_groups = normalize_bullet_groups(
            slide_data.get("bulletGroups") or slide_data.get("bullet_groups") or []
        )
        image_prompts = parse_image_prompts(slide_data)
        if not bullet_groups:
            bullets = expand_bullets(bullets)
        if not bullets and slide_data.get("rawText"):
            raw_text = str(slide_data.get("rawText"))
            bullets = [line.strip() for line in raw_text.splitlines() if line.strip()]

        image_sources = resolve_image_sources(slide_data, temp_files)

        default_layout_index = (
            layout_profiles[min(1, len(layout_profiles) - 1)]["index"]
            if layout_profiles
            else min(1, len(layouts) - 1)
        )
        layout_index = default_layout_index
        if len(layouts) > 0:
            if layout_mode == "template":
                layout_index = fallback_profiles[idx % len(fallback_profiles)]["index"]
            elif layout_mode == "sequential":
                if layout_profiles:
                    layout_index = layout_profiles[idx % len(layout_profiles)]["index"]
                else:
                    layout_index = idx % len(layouts)
            else:
                layout_index = pick_layout_index(
                    slide_data,
                    layout_profiles,
                    default_layout_index,
                )

        if use_existing_slides and idx < existing_count:
            slide = pres.slides[idx]
        else:
            slide = pres.slides.add_slide(layouts[layout_index])

        overflow, title_shape, used_placeholder_tokens = fill_slide_text(
            slide, title, bullets, bullet_groups
        )
        if used_placeholder_tokens:
            _log(f"slide[{idx}] used_text_shapes={len(used_placeholder_tokens)} bullets={len(bullets)} overflow={len(overflow)}")

        if image_prompts:
            notes = slide_data.get("notes", "")
            prompt_text = "配图建议：" + "；".join(image_prompts)
            if notes:
                notes = notes + "\n" + prompt_text
            else:
                notes = prompt_text
            slide_data["notes"] = notes

        if image_sources:
            image_placeholders = collect_image_placeholders(slide, title_shape)
            if image_placeholders:
                for index, image_path in enumerate(image_sources):
                    if index >= len(image_placeholders):
                        break
                    placeholder = image_placeholders[index]
                    try:
                        if hasattr(placeholder, "insert_picture"):
                            placeholder.insert_picture(image_path)
                        else:
                            slide.shapes.add_picture(
                                image_path,
                                placeholder.left,
                                placeholder.top,
                                placeholder.width,
                                placeholder.height,
                            )
                    except Exception:
                        try:
                            slide.shapes.add_picture(
                                image_path,
                                placeholder.left,
                                placeholder.top,
                                placeholder.width,
                                placeholder.height,
                            )
                        except Exception:
                            pass
            else:
                # 模板没有图片占位符时，直接在右侧插入一张配图
                for image_path in image_sources:
                    try:
                        if Inches is not None:
                            # 标准 16:9 页面宽 10 英寸、高 5.625 英寸
                            left = Inches(5.5)
                            top = Inches(1.0)
                            width = Inches(3.5)
                            height = Inches(3.5)
                            slide.shapes.add_picture(image_path, left, top, width, height)
                        else:
                            # 回退：不指定大小和位置，交给 pptx 默认处理
                            slide.shapes.add_picture(image_path)
                        break
                    except Exception:
                        continue

        notes = slide_data.get("notes", "")
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        # Overflow content is placed in notes only; continuation slides are not created
        # so that the final slide count strictly matches the requested page count.
        if overflow:
            notes = slide_data.get("notes", "")
            overflow_text = "溢出内容：" + "；".join([str(item) for item in overflow if str(item).strip()])
            slide_data["notes"] = (notes + "\n" + overflow_text).strip() if notes else overflow_text
            try:
                slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
            except Exception:
                pass

        cleared_placeholders, cleared_textboxes = cleanup_unused_text(
            slide, title_shape, used_placeholder_tokens, clean_template_text
        )
        if cleared_placeholders or cleared_textboxes:
            _log(
                f"slide[{idx}] cleared_placeholders={cleared_placeholders} "
                f"cleared_textboxes={cleared_textboxes}"
            )
        _log(f"slide[{idx}] end")

    if existing_count > 0:
        if use_existing_slides:
            if len(slides) < existing_count:
                for index in range(existing_count - 1, len(slides) - 1, -1):
                    if index >= 0:
                        delete_slide(pres, index)
        else:
            for index in range(existing_count - 1, -1, -1):
                delete_slide(pres, index)

    pres.save(output_path)
    for temp_path in temp_files:
        try:
            os.remove(temp_path)
        except Exception:
            pass
    run_placeholder_qa(output_path, strict_qa=strict_qa)
    _log("build_presentation done")


def main():
    parser = argparse.ArgumentParser(description="Fill PPTX template with slide content.")
    parser.add_argument("--template", required=True, help="Path to PPTX template")
    parser.add_argument("--output", required=True, help="Path to output PPTX")
    parser.add_argument("--data-json", required=True, help="Path to JSON payload")
    parser.add_argument("--strict-qa", action="store_true", help="Exit non-zero if placeholder-like text is found in output")
    args = parser.parse_args()

    template_path = Path(args.template)
    output_path = Path(args.output)
    payload_path = Path(args.data_json)

    if not template_path.exists():
        print(f"Template file not found: {template_path}", file=sys.stderr)
        return 2
    if not zipfile.is_zipfile(template_path):
        print(f"Template file is not a valid .pptx: {template_path}", file=sys.stderr)
        return 4
    if not payload_path.exists():
        print("Payload JSON not found", file=sys.stderr)
        return 3

    with payload_path.open("r", encoding="utf-8") as handle:
        payload = json.load(handle)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        build_presentation(str(template_path), str(output_path), payload, strict_qa=getattr(args, "strict_qa", False))
    except Exception as exc:
        print(f"Failed to build presentation: {exc}", file=sys.stderr)
        raise
    return 0


if __name__ == "__main__":
    sys.exit(main())
